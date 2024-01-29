from dotenv import load_dotenv

load_dotenv()  # take environment variables from .env.

from openai import OpenAI
client = OpenAI(base_url="https://api.together.xyz/v1")
name = input("What is your name? ")
lesson = input("what topic are you teaching? ")
ability = input("what ability are the class? ")
length = input("how long is the lesson? ")
prompt = f"""write me a powerpoint presentation for teaching {ability} ability students {lesson} for {length}.

my name is {name}

return the result in the following format: """ + """
```json
{
  "presentationTitle": "title",
  "author": "name",
  "slides": [
    {
      "slideTitle": "slide title",
      "content": [
        "bullet 1",
        "bullet 2",
        "bullet 3"
      ]
    },
}
```
"""
response = client.chat.completions.create(
  model="mistralai/Mixtral-8x7B-Instruct-v0.1",
  messages=[
    {
      "role": "user",
      "content" : prompt
      # "content": "write the text for each slide of a powerpoint to teach a" + length + "lesson on" + lesson + "for" + ability + "ability secondary class"
    }
  ],
  temperature=1,
  max_tokens=256,
  top_p=1,
  frequency_penalty=0,
  presence_penalty=0
)
# print(response)
print(response.choices[0].message.content)

resp = response.choices[0].message.content
resp = resp.split("```json")[1]
resp = resp.split("```")[0]

import json
result = json.loads(resp)
import pprint
pprint.pprint(result)

from pptx import Presentation
from pptx.util import Inches

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1: Title Slide
slide_layout = presentation.slide_layouts[0]
slide_1 = presentation.slides.add_slide(slide_layout)
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Let's Solve Equations!"
subtitle.text = "A Simple Guide for Solving Equations\nYour Name and Date"

# Slide 2: Introduction to Equations
slide_layout = presentation.slide_layouts[1]
slide_2 = presentation.slides.add_slide(slide_layout)
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "Introduction to Equations"
content.text = "- An equation is like a math sentence.\n- It has numbers, letters (variables), and an equal sign.\n- We use equations to find the missing number."

# Repeat similar steps for the rest of the slides...

# Save the PowerPoint presentation
presentation.save("Solving_Equations_Presentation2.pptx")